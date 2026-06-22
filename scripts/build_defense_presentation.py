"""Сборка презентации для защиты курсовой работы.

Презентация собирается «с нуля» из пустой колоды python-pptx: фон-рамка,
логотип ВШЭ, шапка и весь контент рисуются кодом. Никакого внешнего шаблона
не требуется — оформление лежит в наших же ассетах (docs/assets/).

В углу каждого слайда выводится полное название работы (бренд-имя проекта в
тексте слайдов сознательно не упоминается — по замечанию руководителя).

Запуск:  python -m scripts.build_defense_presentation
         (опционально --out docs/Презентация_защита_EventMind.pptx)
"""
from __future__ import annotations

import argparse
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- палитра -----------------------------------------------------------------
NAVY = RGBColor(0x0F, 0x2C, 0x68)
BLUE = RGBColor(0x44, 0x72, 0xC4)
GREY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEA, 0xEF, 0xF7)
GREEN = RGBColor(0xE3, 0xF0, 0xE6)
SAND = RGBColor(0xFD, 0xEF, 0xE3)
FONT = "HSE Sans"

DIAGRAMS = Path("docs/diagrams")
ASSETS = Path("docs/assets")
COVER_BG = ASSETS / "cover_bg.jpg"
CONTENT_BG = ASSETS / "content_bg.jpg"
LOGO_FULL = ASSETS / "logo_full.png"
LOGO_MONO = ASSETS / "logo_mono.png"

# 16:9
SLIDE_W = 13.333
SLIDE_H = 7.5

# Полное название работы — выводится в углу каждого слайда и на титуле.
FULL_TITLE = (
    "Разработка приложения для агрегации информации об IT-событиях и выдачи "
    "персональных рекомендаций на основе модели пользователя"
)
EVENT_LINE = "Защита курсовой, 2026"

# номера слайдов проставляются вторым проходом (общее число известно в конце)
_PAGE_RUNS: list = []


# --- низкоуровневые помощники ------------------------------------------------
def _set_font(run, size=None, bold=None, color=None, italic=None):
    run.font.name = FONT
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def _clear(tf):
    tf.word_wrap = True
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    return p0


def addrun(para, text, size=None, bold=None, color=None, italic=None):
    run = para.add_run()
    run.text = text
    _set_font(run, size, bold, color, italic)
    return run


def add_textbox(slide, left, top, w, h, anchor=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tb


def add_image_fit(slide, path, area_l, area_t, area_w, area_h):
    """Вписать картинку в прямоугольную область, центрируя."""
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(area_w / iw, area_h / ih)
    w, h = iw * scale, ih * scale
    left = area_l + (area_w - w) / 2
    top = area_t + (area_h - h) / 2
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                    width=Inches(w), height=Inches(h))


def add_c4_crop(slide, frac_box, area):
    """Вырезать из общей C4-диаграммы фрагмент с одним модулем и вписать его.

    Кроп считается на лету из 03_c4_containers.png, поэтому всегда совпадает
    с актуальной диаграммой.
    """
    from PIL import Image
    im = Image.open(DIAGRAMS / "03_c4_containers.png")
    w, h = im.size
    left, top, right, bottom = frac_box
    crop = im.crop((int(left * w), int(top * h), int(right * w), int(bottom * h)))
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as fh:
        crop.save(fh.name)
        path = fh.name
    return add_image_fit(slide, path, *area)


# --- каркас слайда (фон, логотип, шапка, заголовок) --------------------------
def add_background(slide, image):
    slide.shapes.add_picture(str(image), Inches(0), Inches(0),
                             width=Inches(SLIDE_W), height=Inches(SLIDE_H))


def add_header(slide, page_no):
    """Шапка контентного слайда: логотип, полное название работы, событие, №."""
    slide.shapes.add_picture(str(LOGO_MONO), Inches(0.52), Inches(0.2),
                             width=Inches(0.52), height=Inches(0.52))

    tb = add_textbox(slide, 1.2, 0.16, 8.4, 0.74, anchor=MSO_ANCHOR.MIDDLE)
    addrun(_clear(tb.text_frame), FULL_TITLE, 8, bold=True, color=NAVY)

    tb = add_textbox(slide, 9.75, 0.16, 2.25, 0.74, anchor=MSO_ANCHOR.MIDDLE)
    p = _clear(tb.text_frame)
    p.alignment = PP_ALIGN.RIGHT
    addrun(p, EVENT_LINE, 8, color=GREY)

    tb = add_textbox(slide, 12.05, 0.16, 0.85, 0.74, anchor=MSO_ANCHOR.MIDDLE)
    p = _clear(tb.text_frame)
    p.alignment = PP_ALIGN.RIGHT
    _PAGE_RUNS.append(addrun(p, f"{page_no} / {{T}}", 9.5, bold=True, color=NAVY))


def add_heading(slide, text):
    tb = add_textbox(slide, 0.55, 1.0, 12.3, 0.85, anchor=MSO_ANCHOR.MIDDLE)
    addrun(_clear(tb.text_frame), text, 30, bold=True, color=NAVY)


def new_slide(prs, heading=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CONTENT_BG)
    add_header(slide, len(prs.slides))
    if heading:
        add_heading(slide, heading)
    return slide


# --- блоки контента ----------------------------------------------------------
def bullets(slide, left, top, w, h, items, size=16, gap=8, anchor=MSO_ANCHOR.TOP):
    tb = add_textbox(slide, left, top, w, h, anchor=anchor)
    tf = tb.text_frame
    first = True
    for it in items:
        text, level = it if isinstance(it, tuple) else (it, 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        marker = "•  " if level == 0 else "–  "
        addrun(p, marker + text, size if level == 0 else size - 2,
               bold=False, color=NAVY if level == 0 else GREY)
    return tb


def card(slide, left, top, w, h, title, body, fill=LIGHT, title_color=NAVY,
         title_size=15, body_size=12.5):
    sh = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = NAVY
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", Pt(8))
    addrun(_clear(tf), title, title_size, bold=True, color=title_color)
    for line in body:
        bp = tf.add_paragraph()
        bp.space_before = Pt(3)
        sub = line.startswith("– ")
        addrun(bp, line, body_size - (1 if sub else 0), color=GREY)
    return sh


def caption(slide, text, top=7.0):
    tb = add_textbox(slide, 0.55, top, 12.2, 0.4)
    p = _clear(tb.text_frame)
    p.alignment = PP_ALIGN.CENTER
    addrun(p, text, 12, italic=True, color=GREY)


def style_table(table, col_widths=None, header_size=12.5, body_size=11,
                first_col_bold=True):
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    n_cols = len(table.columns)
    for ci in range(n_cols):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                _set_font(r, header_size, bold=True, color=WHITE)
    for ri in range(1, len(table.rows)):
        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _set_font(r, body_size,
                              bold=(first_col_bold and ci == 0),
                              color=NAVY if (first_col_bold and ci == 0) else GREY)


def add_table(slide, data, left, top, w, h, col_widths=None, header_size=12.5,
              body_size=11, highlight_last=False):
    rows, cols = len(data), len(data[0])
    gf = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                Inches(w), Inches(h))
    table = gf.table
    tblPr = table._tbl.tblPr
    tblPr.set("firstRow", "1")
    tblPr.set("bandRow", "0")
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            for m in ("left", "right"):
                setattr(cell.text_frame, f"margin_{m}", Pt(6))
            for m in ("top", "bottom"):
                setattr(cell.text_frame, f"margin_{m}", Pt(3))
    style_table(table, col_widths=col_widths, header_size=header_size,
                body_size=body_size)
    if highlight_last:
        last = len(data) - 1
        for ci in range(cols):
            cell = table.cell(last, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = GREEN
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _set_font(r, body_size, bold=True, color=NAVY)
    return table


# --- титульный слайд ---------------------------------------------------------
def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COVER_BG)

    slide.shapes.add_picture(str(LOGO_FULL), Inches(1.0), Inches(1.05),
                             width=Inches(1.0), height=Inches(1.0))

    tb = add_textbox(slide, 2.25, 1.05, 4.6, 1.0, anchor=MSO_ANCHOR.MIDDLE)
    tf = tb.text_frame
    addrun(_clear(tf), "Национальный исследовательский университет", 11, color=GREY)
    addrun(tf.add_paragraph(), "«Высшая школа экономики» — Пермь", 12,
           bold=True, color=NAVY)

    tb = add_textbox(slide, 7.1, 1.05, 3.0, 1.0, anchor=MSO_ANCHOR.MIDDLE)
    tf = tb.text_frame
    addrun(_clear(tf), "Магистерская программа", 11, color=GREY)
    addrun(tf.add_paragraph(), "«Бизнес-аналитика»", 12, bold=True, color=NAVY)

    tb = add_textbox(slide, 10.6, 1.05, 1.7, 1.0, anchor=MSO_ANCHOR.MIDDLE)
    p = _clear(tb.text_frame)
    p.alignment = PP_ALIGN.RIGHT
    addrun(p, "Пермь, 2026", 12, color=GREY)

    tb = add_textbox(slide, 1.0, 2.55, 11.3, 1.9, anchor=MSO_ANCHOR.MIDDLE)
    addrun(_clear(tb.text_frame), FULL_TITLE, 27, bold=True, color=NAVY)

    tb = add_textbox(slide, 1.0, 5.0, 5.6, 1.0, anchor=MSO_ANCHOR.TOP)
    tf = tb.text_frame
    addrun(_clear(tf), "Работу выполнил студент:", 15, color=GREY)
    addrun(tf.add_paragraph(), "Белов Егор Александрович", 16, bold=True, color=NAVY)

    tb = add_textbox(slide, 7.0, 5.0, 5.3, 1.2, anchor=MSO_ANCHOR.TOP)
    tf = tb.text_frame
    addrun(_clear(tf), "Научный руководитель:", 15, color=GREY)
    addrun(tf.add_paragraph(),
           "Доцент кафедры информационных технологий в бизнесе", 12, color=GREY)
    addrun(tf.add_paragraph(), "Городилов А. Ю.", 16, bold=True, color=NAVY)
    return slide


def build_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, CONTENT_BG)
    add_header(slide, len(prs.slides))

    box = slide.shapes.add_shape(1, Inches(1.7), Inches(2.5), Inches(9.9), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = _clear(tf)
    p.alignment = PP_ALIGN.CENTER
    addrun(p, "Благодарю за внимание!", 40, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    addrun(p2, "Готов ответить на ваши вопросы", 20,
           color=RGBColor(0xCF, 0xDA, 0xEE))

    tb = add_textbox(slide, 1.7, 5.1, 9.9, 1.3)
    tf = tb.text_frame
    p = _clear(tf)
    p.alignment = PP_ALIGN.CENTER
    addrun(p, "Белов Егор Александрович", 16, bold=True, color=NAVY)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    addrun(p2, "НИУ ВШЭ — Пермь · «Бизнес-аналитика»", 14, color=GREY)
    return slide


# ===========================================================================
def build(out: str):
    _PAGE_RUNS.clear()
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_cover(prs)

    # --- Слайд: Актуальность -------------------------------------------------
    s = new_slide(prs, "Актуальность")
    bullets(s, 0.6, 1.95, 7.2, 4.8, [
        "IT-мероприятий очень много: конференции, митапы, вебинары и "
        "хакатоны идут десятками в неделю.",
        "Анонсы разбросаны по множеству сайтов, агрегаторов и каналов — "
        "единой точки входа нет.",
        "Отобрать подходящее по теме, уровню, формату и городу — ручная "
        "рутина, отнимающая время.",
        "Обычные афиши не знают карьерной цели пользователя и его истории "
        "интересов.",
        "Растёт запрос на персонализацию: лента «под меня», а не общий список.",
    ], size=16, gap=12)
    card(s, 8.2, 2.1, 4.6, 2.0, "Проблема",
         ["Подходящие IT-события теряются в потоке анонсов из множества "
          "источников, а их отбор остаётся ручным."], body_size=14)
    card(s, 8.2, 4.45, 4.6, 2.3, "Решение",
         ["Автоматический сбор событий и персональная лента с понятным "
          "объяснением выбора, учитывающая интересы и цель пользователя."],
         fill=GREEN, body_size=14)

    # --- Слайд: Объект и предмет --------------------------------------------
    s = new_slide(prs, "Объект и предмет исследования")
    card(s, 0.7, 2.2, 5.8, 3.0, "Объект исследования",
         ["Процессы агрегации информации об IT-мероприятиях и построения "
          "персонализированных рекомендаций для пользователей."],
         body_size=14)
    card(s, 6.9, 2.2, 5.8, 3.0, "Предмет исследования",
         ["Программная система, объединяющая сбор и нормализацию событий, "
          "гибридный рекомендатель и пользовательский интерфейс вокруг "
          "общей базы данных."], fill=GREEN, body_size=14)

    # --- Слайд: Цель и задачи ------------------------------------------------
    s = new_slide(prs, "Цель и задачи")
    card(s, 0.6, 1.95, 5.5, 3.4, "Цель",
         ["Разработать систему, которая автоматически собирает анонсы "
          "IT-мероприятий из множества источников, нормализует их и строит "
          "для каждого пользователя персональную ленту рекомендаций с "
          "понятным объяснением выбора."], body_size=14)
    tb = add_textbox(s, 6.4, 1.0, 4, 0.6)
    addrun(_clear(tb.text_frame), "Задачи", 20, bold=True, color=NAVY)
    bullets(s, 6.4, 1.7, 6.4, 5.4, [
        "Анализ предметной области и существующих решений",
        "Анализ подходов к построению рекомендаций",
        "Проектирование архитектуры системы",
        "Сбор и автоматическая нормализация событий из разных источников",
        "Реализация гибридного алгоритма рекомендаций",
        "Разработка пользовательского интерфейса и обратной связи",
        "Тестирование и оценка качества рекомендаций",
    ], size=15.5, gap=11)

    # --- Слайд: Проблемы (овалы) --------------------------------------------
    s = new_slide(prs, "Проблемы предметной области")

    def oval(left, top, w, h, text, fill, size=14, tcolor=WHITE, bold=False):
        o = s.shapes.add_shape(9, Inches(left), Inches(top), Inches(w), Inches(h))
        o.fill.solid()
        o.fill.fore_color.rgb = fill
        o.line.color.rgb = WHITE
        o.line.width = Pt(1.25)
        o.shadow.inherit = False
        tf = o.text_frame
        tf.word_wrap = True
        p = _clear(tf)
        p.alignment = PP_ALIGN.CENTER
        addrun(p, text, size, bold=bold, color=tcolor)
        return o

    oval(4.42, 3.55, 4.5, 1.5, "Проблемы", NAVY, size=27, bold=True)
    oval(0.5, 2.0, 3.6, 1.45, "Разрозненность источников событий", BLUE)
    oval(0.65, 5.2, 3.6, 1.45, "Информационный шум и дубли анонсов", BLUE)
    oval(4.9, 5.95, 3.55, 1.4, "Нет учёта уровня и карьерной цели", BLUE)
    oval(9.2, 2.0, 3.6, 1.45, "Рекомендации без объяснения выбора", BLUE)
    oval(9.05, 5.2, 3.6, 1.45, "Ручная фильтрация отнимает время", BLUE)

    # --- Слайд: Анализ существующих решений ---------------------------------
    s = new_slide(prs, "Анализ существующих решений")
    data = [
        ["Решение", "Сильные стороны", "Ограничения"],
        ["Универсальные площадки\n(Timepad, Meetup, Eventbrite)",
         "Фильтры по категориям, городам и датам",
         "Не специализированы на IT, не учитывают поведение пользователя"],
        ["Тематические IT-площадки\n(Habr Events)",
         "Релевантный контент по теме",
         "Нет модели интересов и регулярных персональных подборок"],
        ["Профессиональные сети\n(LinkedIn Events)",
         "Связь с профессиональным профилем",
         "Фокус на карьерных событиях, мало митапов и конференций"],
        ["Каналы сообществ\n(Telegram)",
         "Оперативные анонсы из первых рук",
         "Нет организованного поиска и персонализации"],
        ["Разрабатываемая система",
         "Агрегация из источников, учёт реакций, объяснимые рекомендации",
         "Снимает разрозненность, добавляет персонализацию и дайджест"],
    ]
    add_table(s, data, 0.55, 1.95, 12.25, 4.6,
              col_widths=[3.2, 4.35, 4.7], body_size=11.5, highlight_last=True)
    caption(s, "Ни одно из решений не объединяет агрегацию из разных "
               "источников, учёт неявной обратной связи и объяснимость.")

    # --- Слайд: Анализ подходов ---------------------------------------------
    s = new_slide(prs, "Анализ подходов к рекомендациям")
    data = [
        ["Подход", "Преимущества", "Недостатки", "В системе"],
        ["По содержанию\n(rule + cosine)", "Работает при малом числе "
         "пользователей, объясним", "Ограниченное разнообразие",
         "Базовые правила + близость"],
        ["Коллаборативная\nфильтрация", "Скрытые закономерности, "
         "неожиданные находки", "Холодный старт, нужны объёмы",
         "LightGCN (эксперим.)"],
        ["Семантический\n(эмбеддинги)", "Учёт смысла, многоязычность",
         "Затраты на векторы", "MiniLM + pgvector"],
        ["Контекстный\nбандит", "Баланс новизны и точности, онлайн-адаптация",
         "Настройка признаков", "LinUCB поверх гибрида"],
        ["Долговременная\nпамять (mem0)", "Запоминает цели и интересы",
         "Затраты на LLM, риск устаревания", "Семантический поиск"],
        ["Гибридный\nподход", "Точность, устойчивость к холодному старту",
         "Сложность настройки весов", "Основной алгоритм"],
    ]
    add_table(s, data, 0.55, 1.9, 12.25, 5.0,
              col_widths=[2.3, 3.9, 2.95, 3.1], body_size=10.5,
              highlight_last=True)

    # --- Слайд: Use case -----------------------------------------------------
    s = new_slide(prs, "Диаграмма вариантов использования")
    add_image_fit(s, DIAGRAMS / "01_use_case.png", 0.5, 1.8, 7.6, 5.2)
    card(s, 8.4, 2.0, 4.4, 2.3, "Пользователь",
         ["Настройка профиля, персональная лента, объяснения, оценки и "
          "сохранения, поиск обычным языком, дайджест."], body_size=13)
    card(s, 8.4, 4.5, 4.4, 1.9, "Администратор",
         ["Загрузка событий из источников и просмотр аналитики по системе."],
         fill=GREEN, body_size=13)

    # --- Слайд: Бизнес-процесс ----------------------------------------------
    s = new_slide(prs, "Бизнес-процесс системы")
    add_image_fit(s, DIAGRAMS / "02_business_process.png", 0.6, 1.85, 12.1, 5.0)
    caption(s, "От сбора и нормализации событий до доставки персональной "
               "ленты и обработки обратной связи.")

    # --- Слайды: Архитектура (общая схема слева + приближённый модуль) -------
    def arch_slide(title, frac_box, module_name):
        s = new_slide(prs, title)
        add_image_fit(s, DIAGRAMS / "03_c4_containers.png", 0.3, 1.7, 6.0, 5.3)
        add_c4_crop(s, frac_box, (6.7, 1.85, 6.3, 4.95))
        caption(s, f"Слева — общая схема на уровне контейнеров; справа — "
                   f"приближён модуль «{module_name}».")
        return s

    arch_slide("Архитектура: серверная часть",
               (0.09, 0.51, 0.47, 0.77), "серверная часть")
    arch_slide("Архитектура: Telegram-бот",
               (0.02, 0.333, 0.285, 0.52), "Telegram-бот")
    arch_slide("Архитектура: планировщик",
               (0.315, 0.347, 0.625, 0.535), "планировщик")

    # --- Слайд: Модель данных (ER) ------------------------------------------
    s = new_slide(prs, "Модель данных")
    add_image_fit(s, DIAGRAMS / "04_er_diagram.png", 0.5, 1.9, 12.3, 4.9)
    caption(s, "Связи проведены от ключа к ключу (PK → FK): пользователи, "
               "события, взаимодействия, векторы, состояние бандита и память.")

    # --- Слайд: 9 компонент рекомендатора -----------------------------------
    s = new_slide(prs, "Гибридный рекомендатель: 9 компонент")
    data = [
        ["Компонент", "Парадигма / источник", "Назначение"],
        ["rule", "по правилам, профиль", "Совпадение явных предпочтений"],
        ["cosine", "по содержанию, эмбеддинги", "Семантическая близость"],
        ["bayesian", "Thompson sampling", "Статистика по темам с затуханием"],
        ["quality", "оценка LLM 1–10", "Качество описания и источника"],
        ["hype", "оценка LLM 1–10", "Актуальность темы сейчас"],
        ["freshness", "затухание по дате", "Свежесть события"],
        ["skill_gap", "учёт карьерной цели", "Соответствие цели пользователя"],
        ["bandit", "контекстный LinUCB", "Баланс новизны и точности"],
        ["gnn", "коллаборативный LightGCN", "Коллаборативный сигнал (эксперим.)"],
    ]
    add_table(s, data, 0.55, 1.85, 9.0, 5.15,
              col_widths=[1.9, 3.5, 3.6], body_size=11)
    bullets(s, 9.8, 2.0, 3.2, 5.0, [
        "Веса задаются конфигурацией",
        "Переранжирование MMR (λ=0.7) для разнообразия",
        "Антиповтор серий: один выпуск из серии",
        "Каждая компонента изолирована — сбой не ломает выдачу",
        "Кэш выдачи на 15 минут",
    ], size=12.5, gap=10)

    # --- Слайд: Sequence рекомендации ---------------------------------------
    s = new_slide(prs, "Как строится рекомендация")
    add_image_fit(s, DIAGRAMS / "05_sequence_recommend.png", 0.4, 1.8, 12.5, 5.1)
    caption(s, "Запрос ленты → гибрид отбирает top-N → один вызов LLM пишет "
               "объяснения пакетом (по найденным фактам, без выдумок).")

    # --- Слайд: Ingestion ----------------------------------------------------
    s = new_slide(prs, "Сбор и нормализация событий")
    add_image_fit(s, DIAGRAMS / "06_activity_normalization.png", 0.4, 1.95, 12.5, 2.9)
    bullets(s, 0.9, 5.1, 5.9, 1.9, [
        "6 источников: habr, rss, kudago, luma, meetup, telegram",
        "Сырьё → таблица raw_events → нормализация через LLM",
        "Проверка типов и дат (Pydantic)",
    ], size=13, gap=7)
    bullets(s, 7.0, 5.1, 5.6, 1.9, [
        "Поиск дубликатов по вектору (pgvector, cosine ≥ 0.92)",
        "Идемпотентный повторный сбор",
        "Пакетная нормализация с адаптивным размером порции",
    ], size=13, gap=7)

    # --- Слайд: LLM-цепочка --------------------------------------------------
    s = new_slide(prs, "Цепочка LLM и надёжность")
    card(s, 0.6, 1.95, 6.0, 3.7, "Цепочка провайдеров",
         ["1.  Gemini — основной (транспорт REST, автоподбор модели)",
          "2.  Groq llama-3.3-70b — резерв",
          "3.  Groq llama-3.1-8b — крайний резерв",
          "",
          "Предохранитель: 5 сбоёв подряд → пауза 120 с.",
          "Пауза по провайдеру: 2 сбоя → пропуск на 10 минут."],
         body_size=13.5)
    bullets(s, 6.9, 2.0, 6.0, 4.8, [
        ("Плавная деградация на двух уровнях:", 0),
        ("LLM → следующий провайдер цепочки", 1),
        ("гибрид → выдача только по правилам при сбое", 1),
        ("Общий секрет X-API-Key + логирование запросов", 0),
        ("Очистка пользовательского ввода от инъекций в промпт", 0),
        ("Выдача только на чтение + кэш рекомендаций", 0),
    ], size=14, gap=10)

    # --- Слайд: Оценка качества ---------------------------------------------
    s = new_slide(prs, "Оценка качества")
    tb = add_textbox(s, 0.6, 1.65, 8, 0.4)
    addrun(_clear(tb.text_frame), "Оффлайн-оценка (Recall@k / nDCG@k)",
           14, bold=True, color=NAVY)
    t1 = [
        ["Алгоритм", "Recall@5", "Recall@10", "nDCG@5", "nDCG@10"],
        ["rule", "0,50", "0,80", "0,26", "0,36"],
        ["hybrid", "0,65", "0,85", "0,38", "0,44"],
        ["bayesian", "0,65", "0,85", "0,38", "0,45"],
    ]
    add_table(s, t1, 0.6, 2.1, 12.2, 1.9,
              col_widths=[2.4, 2.45, 2.45, 2.45, 2.45], body_size=12.5)

    tb = add_textbox(s, 0.6, 4.25, 8, 0.4)
    addrun(_clear(tb.text_frame), "Оценка LLM-судьёй (шкала 1–5)",
           14, bold=True, color=NAVY)
    t2 = [
        ["Алгоритм", "Релевантность", "Разнообразие"],
        ["rule", "4,33", "2,53"],
        ["hybrid", "4,40", "2,10"],
        ["bayesian", "4,40", "2,03"],
    ]
    add_table(s, t2, 0.6, 4.7, 7.0, 1.7,
              col_widths=[2.4, 2.3, 2.3], body_size=12.5)
    bullets(s, 7.9, 4.7, 4.9, 2.0, [
        "Гибрид превосходит базовое правило по Recall@5 (0,65 против 0,50) "
        "и nDCG@5 (0,38 против 0,26).",
        "369 автотестов; оценка по воспроизводимым скриптам (seed=42).",
    ], size=12.5, gap=9)

    # --- Слайд: Заключение ---------------------------------------------------
    s = new_slide(prs, "Заключение")
    bullets(s, 0.6, 1.95, 6.1, 5.0, [
        "Все поставленные задачи выполнены.",
        "Собрана рабочая система: сервер, бот и планировщик вокруг общей "
        "базы данных (PostgreSQL + pgvector).",
        "Реализован гибридный рекомендатель из 9 компонент с объяснением "
        "выбора.",
        "Сбор из 6 источников с автоматической нормализацией и дедупликацией.",
        "Качество подтверждено метриками и 369 тестами.",
    ], size=15.5, gap=12)
    card(s, 7.0, 1.95, 5.8, 3.6, "Направления развития",
         ["• Расширение списка источников событий",
          "• Включение коллаборативного сигнала на большем наборе данных",
          "• Расширение пользовательских сценариев бота",
          "• Сравнительное тестирование весов рекомендатора",
          "• Развитие в магистерскую работу"],
         fill=SAND, body_size=13.5)

    # --- Слайд: Спасибо ------------------------------------------------------
    build_thanks(prs)

    # проставить общее число слайдов
    total = len(prs.slides)
    for run in _PAGE_RUNS:
        run.text = run.text.replace("{T}", str(total))

    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"saved: {out}  ({total} слайдов)")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="docs/Презентация_защита_EventMind.pptx")
    args = ap.parse_args()
    build(args.out)
